import asyncio
import base64
import shutil
import subprocess
import tempfile
from abc import ABC, abstractmethod
from io import BytesIO
from pathlib import Path
from typing import Any, cast

import fitz  # type: ignore[import-untyped]
import pytesseract  # type: ignore[import-untyped]
from openai import AsyncOpenAI
from PIL import Image
from pptx import Presentation
from pydantic import BaseModel, Field

from app.intelligence.models import ExtractedSection


class DocumentExtractionError(ValueError):
    pass


class VisualAnalyzer(ABC):
    @abstractmethod
    async def analyze(self, image: bytes, native_text: str) -> str: ...


class LocalOcrVisualAnalyzer(VisualAnalyzer):
    async def analyze(self, image: bytes, native_text: str) -> str:
        def _ocr() -> str:
            with Image.open(BytesIO(image)) as page:
                return str(pytesseract.image_to_string(page)).strip()

        ocr = await asyncio.to_thread(_ocr)
        if not ocr:
            return "Visual page with no additional readable labels."
        if native_text and _normalized(ocr) in _normalized(native_text):
            return "Visual content accompanies the extracted page text."
        return f"Visual text: {ocr[:3000]}"


class _VisionSummary(BaseModel):
    summary: str = Field(min_length=1, max_length=3000)


class OpenAIVisualAnalyzer(VisualAnalyzer):
    def __init__(self, api_key: str, model: str) -> None:
        self._client = AsyncOpenAI(api_key=api_key)
        self._model = model

    async def analyze(self, image: bytes, native_text: str) -> str:
        encoded = base64.b64encode(image).decode("ascii")
        response = await self._client.responses.parse(
            model=self._model,
            input=cast(
                Any,
                [
                    {
                        "role": "system",
                        "content": (
                            "Describe the pedagogically meaningful visual content on this lecture "
                            "page or slide. Capture diagram relationships, chart meaning, "
                            "formulas, labels, and examples that native text extraction misses. "
                            "Do not infer facts that are not visible. Return one concise "
                            "source-grounding summary."
                        ),
                    },
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "input_text",
                                "text": f"Native extracted text:\n{native_text[:8000]}",
                            },
                            {
                                "type": "input_image",
                                "image_url": f"data:image/png;base64,{encoded}",
                            },
                        ],
                    },
                ],
            ),
            text_format=_VisionSummary,
        )
        parsed = response.output_parsed
        if parsed is None:
            raise RuntimeError("Visual analysis did not match the expected schema.")
        return parsed.summary


class DocumentExtractor:
    def __init__(self, analyzer: VisualAnalyzer, max_pages: int = 200) -> None:
        self._analyzer = analyzer
        self._max_pages = max_pages

    async def extract(self, path: Path, source_type: str) -> tuple[ExtractedSection, ...]:
        if source_type == "pdf":
            pages = await asyncio.to_thread(_read_pdf, path, self._max_pages)
        elif source_type == "pptx":
            pages = await asyncio.to_thread(_read_pptx, path, self._max_pages)
        else:
            raise DocumentExtractionError(f"Unsupported supplemental source type: {source_type}")

        sections: list[ExtractedSection] = []
        for index, page in enumerate(pages):
            visual_summary = ""
            if page.image and (len(page.native_text.strip()) < 80 or page.has_visuals):
                visual_summary = await self._analyzer.analyze(page.image, page.native_text)
            sections.append(
                ExtractedSection(
                    section_index=index,
                    page_number=index + 1,
                    title=page.title,
                    native_text=page.native_text,
                    speaker_notes=page.speaker_notes,
                    visual_summary=visual_summary,
                    metadata={"has_visuals": page.has_visuals},
                )
            )
        return tuple(sections)


class _Page:
    def __init__(
        self,
        *,
        title: str | None,
        native_text: str,
        speaker_notes: str = "",
        image: bytes | None = None,
        has_visuals: bool = False,
    ) -> None:
        self.title = title
        self.native_text = native_text
        self.speaker_notes = speaker_notes
        self.image = image
        self.has_visuals = has_visuals


def _read_pdf(path: Path, max_pages: int) -> list[_Page]:
    try:
        document = fitz.open(path)
    except Exception as exc:
        raise DocumentExtractionError("The PDF could not be opened.") from exc
    with document:
        if document.needs_pass:
            raise DocumentExtractionError("Encrypted PDFs are not supported.")
        if document.page_count > max_pages:
            raise DocumentExtractionError(f"Documents may contain at most {max_pages} pages.")
        pages: list[_Page] = []
        for page in document:
            text = page.get_text("text").strip()
            title = next((line.strip() for line in text.splitlines() if line.strip()), None)
            pixmap = page.get_pixmap(matrix=fitz.Matrix(1.35, 1.35), alpha=False)
            pages.append(
                _Page(
                    title=title[:180] if title else None,
                    native_text=text,
                    image=pixmap.tobytes("png"),
                    has_visuals=bool(page.get_images(full=True) or page.get_drawings()),
                )
            )
        return pages


def _read_pptx(path: Path, max_pages: int) -> list[_Page]:
    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        raise DocumentExtractionError("The PowerPoint file could not be opened.") from exc
    if len(presentation.slides) > max_pages:
        raise DocumentExtractionError(f"Documents may contain at most {max_pages} slides.")

    rendered = _render_pptx(path)
    pages: list[_Page] = []
    for index, slide in enumerate(presentation.slides):
        chunks: list[str] = []
        title: str | None = None
        has_visuals = False
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                value = str(shape.text).strip()
                if value:
                    chunks.append(value)
                    if title is None:
                        title = value.splitlines()[0][:180]
            if getattr(shape, "shape_type", None) in {3, 6, 13, 19}:
                has_visuals = True
        notes = ""
        try:
            notes = str(slide.notes_slide.notes_text_frame.text).strip()
        except (AttributeError, ValueError):
            notes = ""
        rendered_page = rendered[index] if index < len(rendered) else None
        pages.append(
            _Page(
                title=title,
                native_text="\n".join(chunks),
                speaker_notes=notes,
                image=rendered_page.image if rendered_page else None,
                has_visuals=has_visuals or bool(rendered_page and rendered_page.has_visuals),
            )
        )
    return pages


def _render_pptx(path: Path) -> list[_Page]:
    executable = shutil.which("soffice") or shutil.which("libreoffice")
    if executable is None:
        return []
    with tempfile.TemporaryDirectory(prefix="manifold-slides-") as directory:
        result = subprocess.run(
            [executable, "--headless", "--convert-to", "pdf", "--outdir", directory, str(path)],
            check=False,
            capture_output=True,
            text=True,
            timeout=120,
        )
        output = Path(directory) / f"{path.stem}.pdf"
        if result.returncode != 0 or not output.exists():
            return []
        return _read_pdf(output, 10_000)


def _normalized(value: str) -> str:
    return " ".join(value.lower().split())
