from pathlib import Path
from uuid import uuid4

import fitz
import pytest
from pptx import Presentation
from pptx.util import Inches

from app.dashboard.models import (
    DashboardSignal,
    DashboardSignalStatus,
    DashboardSignalType,
    TopicHealth,
)
from app.dashboard.priority_generation import generate_priorities
from app.intelligence.extractor import DocumentExtractionError, DocumentExtractor, VisualAnalyzer
from app.intelligence.local_agent import LocalCourseImprovementAgent
from app.intelligence.openai_agent import _ImprovementOutput


class RecordingAnalyzer(VisualAnalyzer):
    def __init__(self) -> None:
        self.calls = 0

    async def analyze(self, image: bytes, native_text: str) -> str:
        self.calls += 1
        return f"Observed visual beside: {native_text[:30]}"


def test_openai_improvement_schema_is_strict_and_uses_json_object_string() -> None:
    schema = _ImprovementOutput.model_json_schema()

    assert schema["additionalProperties"] is False
    assert schema["properties"]["proposed_state_json"]["type"] == "string"


@pytest.mark.anyio
async def test_pdf_and_powerpoint_sources_preserve_page_and_slide_grounding(
    tmp_path: Path,
) -> None:
    pdf_path = tmp_path / "context.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "Confidence and retrieval practice evidence")
    document.save(pdf_path)
    document.close()

    deck_path = tmp_path / "lecture.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1))
    box.text = "A worked example of deliberate practice"
    deck.save(deck_path)

    analyzer = RecordingAnalyzer()
    extractor = DocumentExtractor(analyzer, max_pages=10)
    pdf_sections = await extractor.extract(pdf_path, "pdf")
    deck_sections = await extractor.extract(deck_path, "pptx")

    assert pdf_sections[0].page_number == 1
    assert "retrieval practice" in pdf_sections[0].native_text
    assert deck_sections[0].page_number == 1
    assert "worked example" in deck_sections[0].native_text
    assert analyzer.calls >= 1


@pytest.mark.anyio
async def test_document_page_limit_is_enforced(tmp_path: Path) -> None:
    path = tmp_path / "long.pdf"
    document = fitz.open()
    document.new_page()
    document.new_page()
    document.save(path)
    document.close()

    with pytest.raises(DocumentExtractionError, match="at most 1 pages"):
        await DocumentExtractor(RecordingAnalyzer(), max_pages=1).extract(path, "pdf")


def test_priority_brief_combines_persisted_signals_and_topic_learning_risk() -> None:
    course_id = uuid4()
    topic_id = uuid4()
    logical_id = uuid4()
    signal = DashboardSignal(
        id=uuid4(),
        course_id=course_id,
        type=DashboardSignalType.UNDERPERFORMING_CONTENT,
        related_entity_type="question",
        related_entity_id=uuid4(),
        status=DashboardSignalStatus.OPEN,
        ai_diagnosis={
            "title": "Question is producing uncertainty",
            "summary": "Learner evidence crossed the review threshold.",
            "metrics": {"attempts": 8, "struggling_learners": 2},
            "target_logical_artifact_id": str(logical_id),
        },
        instructor_action=None,
    )
    topic = TopicHealth(
        topic_id=topic_id,
        logical_id=logical_id,
        title="Deliberate practice",
        learner_reach=3,
        attempts=10,
        correct_attempts=4,
        confidence_1=2,
        confidence_2=3,
        confidence_3=3,
        confidence_4=2,
        mastered_learners=1,
        practiced_learners=1,
        struggling_learners=2,
        remediation_attempts=4,
        active_clips=1,
        clip_duration_seconds=90,
        assessment_count=1,
        concept_count=2,
    )

    priorities = generate_priorities((signal,), (topic,))

    assert priorities[0].specialist_role == "assessment_designer"
    assert priorities[0].target_logical_artifact_id == logical_id
    assert any(item.target_logical_artifact_id == logical_id for item in priorities)
    assert any(item.affected_learners == 3 for item in priorities)


@pytest.mark.anyio
async def test_local_specialist_prepares_reversible_typed_topic_diff() -> None:
    logical_id = uuid4()
    draft = await LocalCourseImprovementAgent().prepare(
        artifact_type="topic",
        logical_artifact_id=logical_id,
        current_state={"summary": "Learners compare practice strategies."},
        evidence={"incorrect_attempts": 4},
        instruction="Clarify the difference between repetition and deliberate practice.",
    )

    assert draft.logical_artifact_id == logical_id
    assert draft.before_state["summary"] == "Learners compare practice strategies."
    assert "Teaching focus" in str(draft.proposed_state["summary"])
    assert draft.before_state != draft.proposed_state
